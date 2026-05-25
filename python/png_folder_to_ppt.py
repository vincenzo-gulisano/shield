#!/usr/bin/env python3

"""
Build a PowerPoint deck from PNG images in a directory tree.

The script recursively scans an input directory, finds every `.png` file, and
creates one slide per image. Images are ordered by folder path, and within each
folder they are sorted by file name. Each image is scaled as large as possible
while preserving its aspect ratio and leaving a small caption area at the bottom
of the slide. The caption contains the image path relative to the scanned input
directory, so the slide can be traced back to the original file.

Usage:
    python3 python/png_folder_to_ppt.py <input-directory> <output-pptx>

Dependency:
    pip install python-pptx
"""

import argparse
from pathlib import Path


def find_pngs(input_dir: Path) -> list[Path]:
    """Return PNG files recursively, sorted by folder path and then file name."""
    png_paths = [path for path in input_dir.rglob("*.png") if path.is_file()]

    # Sort by the parent folder first, then by the file name inside that folder.
    # This preserves a predictable directory traversal while matching the user's
    # requested ordering within each folder.
    return sorted(png_paths, key=lambda path: (path.parent.as_posix().lower(), path.name.lower()))


def fit_size(image_width: int, image_height: int, max_width: int, max_height: int) -> tuple[int, int]:
    """Compute the largest size that fits inside max_width/max_height without distortion."""
    image_aspect = image_width / image_height
    box_aspect = max_width / max_height

    # If the image is proportionally wider than the box, width is limiting.
    # Otherwise, height is limiting.
    if image_aspect > box_aspect:
        width = max_width
        height = round(max_width / image_aspect)
    else:
        height = max_height
        width = round(max_height * image_aspect)
    return width, height


def add_image_slide(prs, png_path: Path, caption: str) -> None:
    """Add one blank slide containing a fitted PNG and its source path caption."""
    from PIL import Image
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Reserve a thin strip for the path caption below the image.
    margin = round(slide_width * 0.025)
    caption_height = round(slide_height * 0.10)
    caption_top = slide_height - caption_height
    image_max_width = slide_width - margin * 2
    image_max_height = caption_top - margin * 2

    with Image.open(png_path) as image:
        image_width_px, image_height_px = image.size

    image_width, image_height = fit_size(
        image_width_px,
        image_height_px,
        image_max_width,
        image_max_height,
    )

    # Center the image in the available image area above the caption.
    image_left = round((slide_width - image_width) / 2)
    image_top = round(margin + (image_max_height - image_height) / 2)
    slide.shapes.add_picture(str(png_path), image_left, image_top, width=image_width, height=image_height)

    # Put the source path below the image. This text is deliberately small so it
    # acts as provenance without stealing much space from the visual.
    text_box = slide.shapes.add_textbox(margin, caption_top, slide_width - margin * 2, caption_height)
    text_frame = text_box.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = caption
    run.font.size = Pt(9)


def build_ppt(input_dir: Path, output_pptx: Path) -> int:
    """Create the PowerPoint deck and return the number of images added."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ModuleNotFoundError as error:
        raise ModuleNotFoundError("Missing dependency. Install it with: pip install python-pptx") from error

    png_paths = find_pngs(input_dir)
    if not png_paths:
        raise ValueError(f"No PNG files found under {input_dir}")

    prs = Presentation()

    # Use a standard widescreen deck. Images are fitted relative to this slide
    # size, while preserving their original aspect ratio.
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    for png_path in png_paths:
        caption = png_path.relative_to(input_dir).as_posix()
        add_image_slide(prs, png_path, caption)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)
    return len(png_paths)


def main() -> None:
    """Parse arguments and run the PNG-to-PowerPoint conversion."""
    parser = argparse.ArgumentParser(
        description="Recursively add PNG files from a directory into a PowerPoint deck."
    )
    parser.add_argument("input_directory", type=Path, help="Directory to scan recursively for PNG files")
    parser.add_argument("output_pptx", type=Path, help="PowerPoint file to create")
    args = parser.parse_args()

    input_dir = args.input_directory.resolve()
    output_pptx = args.output_pptx.resolve()
    if not input_dir.is_dir():
        raise ValueError(f"Input directory does not exist or is not a directory: {input_dir}")

    count = build_ppt(input_dir, output_pptx)
    print(f"Created {output_pptx} with {count} slide(s).")


if __name__ == "__main__":
    main()
